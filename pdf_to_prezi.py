import sys
import subprocess
from pathlib import Path

# Function to install packages
def install_package(package_name):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])

# Try to import necessary packages and install them if missing
try:
    import fitz  # PyMuPDF
except ImportError:
    install_package("pymupdf")
    import fitz

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    install_package("python-pptx")
    from pptx import Presentation
    from pptx.util import Inches

try:
    from PIL import Image
except ImportError:
    install_package("pillow")
    from PIL import Image

try:
    import cv2
    import numpy as np
except ImportError:
    install_package("opencv-python-headless")
    install_package("numpy")
    import cv2
    import numpy as np

# Function to process and crop each page image
def process_pdf_page(page, dpi=300):
    print("  - Rendering page at high resolution...")
    pix = page.get_pixmap(dpi=dpi)
    img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape((pix.height, pix.width, pix.n))
    
    # Convert to grayscale
    gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
    
    # Apply adaptive thresholding for white space detection
    print("  - Applying adaptive thresholding...")
    adaptive_thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                            cv2.THRESH_BINARY_INV, 11, 2)
    
    # Find contours for cropping
    print("  - Finding contours for cropping...")
    contours, _ = cv2.findContours(adaptive_thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    # Determine bounding box for the largest contour
    if contours:
        largest_contour = max(contours, key=cv2.contourArea)
        x, y, w, h = cv2.boundingRect(largest_contour)
        print("  - Cropping the image based on detected content...")
        cropped_img = img_array[y:y+h, x:x+w]
    else:
        print("  - No contours found, using the entire image...")
        cropped_img = img_array

    return Image.fromarray(cropped_img)

# Set paths
pdf_folder = Path('Path/to/your/folder/with/pdf/images')
output_folder = Path('Path/to/save/presentation')
output_folder.mkdir(parents=True, exist_ok=True)

# Create a new PowerPoint presentation
prs = Presentation()

# Process each PDF file in the folder
for pdf_file in pdf_folder.rglob('*.pdf'):
    print(f"Processing PDF file: {pdf_file.name}")
    pdf_document = fitz.open(pdf_file)
    for page_num in range(pdf_document.page_count):
        print(f" - Processing page {page_num + 1} of {pdf_document.page_count}")
        page = pdf_document[page_num]

        # Process and crop the page
        cropped_image = process_pdf_page(page)
        
        # Save the processed image (optional)
        image_path = output_folder / f"{pdf_file.stem}_page_{page_num + 1}.png"
        print(f"  - Saving cropped image to: {image_path}")
        cropped_image.save(image_path)
        
        # Add image to PowerPoint
        print("  - Adding image to PowerPoint slide...")
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        left = Inches(1)
        top = Inches(1)
        slide.shapes.add_picture(str(image_path), left, top, height=Inches(5.5))

    pdf_document.close()
    print(f"Finished processing {pdf_file.name}\n")

# Save the PowerPoint presentation
presentation_path = output_folder / "Cropped_Images.pptx"
prs.save(presentation_path)
print(f"Presentation saved at {presentation_path}")
